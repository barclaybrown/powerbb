import argparse
import json
import logging
import os
from typing import List, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

__version__ = "0.1.0"

log = logging.getLogger(__name__)

EMU_PER_INCH = 914400

def _layout_token(prs: Presentation, slide_layout) -> str:
    for m, master in enumerate(prs.slide_masters):
        for l, layout in enumerate(master.slide_layouts):
            if layout == slide_layout:
                return f"{m}:{l}"
    return "?:?"

def _slide_list(prs: Presentation) -> List[Dict[str, str]]:
    rows = []
    for idx, slide in enumerate(prs.slides, start=1):
        layout_token = _layout_token(prs, slide.slide_layout)
        layout_name = slide.slide_layout.name or ""
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        rows.append({
            "index": idx,
            "layout_index": layout_token,
            "layout_name": layout_name,
            "title": title,
        })
    return rows

def list_slides(prs: Presentation) -> None:
    rows = _slide_list(prs)
    for row in rows:
        line = f"{row['index']:>2}: [{row['layout_index']}] {row['layout_name']}"
        if row['title']:
            line += f" — {row['title']}"
        print(line)
    return rows

def _emu_to_inches(val):
    return float(val) / EMU_PER_INCH

def show_slide(prs: Presentation, number: int) -> None:
    if number < 1 or number > len(prs.slides):
        raise ValueError("Slide number out of range")
    slide = prs.slides[number - 1]
    print(f"Slide {number}: [{_layout_token(prs, slide.slide_layout)}] {slide.slide_layout.name}")
    for shp in slide.shapes:
        shape_type = MSO_SHAPE_TYPE(shp.shape_type).name
        text = shp.text.strip().replace("\n", "\\n") if getattr(shp, 'has_text_frame', False) and shp.text else ""
        pos = (
            _emu_to_inches(shp.left),
            _emu_to_inches(shp.top),
            _emu_to_inches(shp.width),
            _emu_to_inches(shp.height),
        )
        print(f"- {shape_type}: '{text}' at L{pos[0]:.2f} T{pos[1]:.2f} W{pos[2]:.2f} H{pos[3]:.2f}")

def write_json(prs: Presentation, path: str) -> None:
    rows = _slide_list(prs)
    out_dir = os.path.dirname(os.path.abspath(path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(rows, f, indent=2)
    log.info("Wrote %s", os.path.abspath(path))

def interactive_main(prs: Presentation) -> None:
    while True:
        cmd = input("Command (list/show N/json PATH/quit): ").strip()
        if not cmd:
            continue
        if cmd == "list":
            list_slides(prs)
        elif cmd.startswith("show"):
            try:
                _, num = cmd.split()
                show_slide(prs, int(num))
            except Exception:
                print("Usage: show N")
        elif cmd.startswith("json"):
            try:
                _, out = cmd.split(maxsplit=1)
                write_json(prs, out)
            except Exception:
                print("Usage: json OUT.json")
        elif cmd in {"quit", "exit"}:
            break
        else:
            print("Unknown command")

def main(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="PowerPoint slide inspector")
    parser.add_argument("--file", help="Path to PowerPoint file")
    parser.add_argument("--list", action="store_true", help="List slides")
    parser.add_argument("--show", type=int, metavar="N", help="Show slide N")
    parser.add_argument("--json", metavar="OUT.json", help="Write slide list to JSON")
    parser.add_argument("--log-level", default="INFO", help="Logging level")
    parser.add_argument("--version", action="store_true", help="Show version and exit")
    args = parser.parse_args(argv)

    logging.basicConfig(level=getattr(logging, args.log_level.upper(), logging.INFO))

    if args.version:
        print(__version__)
        return 0

    non_interactive = args.list or args.show is not None or args.json
    if non_interactive and not args.file:
        parser.error("--file is required when using --list, --show, or --json")
    if args.file:
        prs = Presentation(os.path.expanduser(args.file))
    else:
        path = input("Path to PowerPoint file: ").strip()
        if not path:
            parser.error("Missing file path")
        prs = Presentation(os.path.expanduser(path))

    if non_interactive:
        if args.list:
            list_slides(prs)
        if args.show is not None:
            show_slide(prs, args.show)
        if args.json:
            write_json(prs, args.json)
    else:
        interactive_main(prs)
    return 0

if __name__ == "__main__":
    raise SystemExit(main())
