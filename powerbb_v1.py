#!/usr/bin/env python3
# powerbb.py

"""

How to use

Generate a conversation prompt (stdout):

python powerbb.py --promptgen --template path/to/YourTemplate.pptx


Write the prompt to a file:

python powerbb.py --promptgen --template path/to/YourTemplate.pptx --prompt-out prompt_for_llm.txt


(Optional) Also emit full template profile JSON (for archival/sharing):

python powerbb.py --dump-layouts --dump-layouts-json template_profile.json --template path/to/YourTemplate.pptx


Paste the produced prompt into your new chat. It explains the powerbb schema and gives the model everything it needs 
from your template (real layout names, alias suggestions, defaults, slide size) so it can generate compatible PowerBB 
JSON on the first try.

Usage examples

Normal (file input):

python powerbb.py --json Q:\path\slides.powerbb.json --output Q:\path\deck.pptx --template Q:\path\ExecDeck.pptx -v


Interactive clipboard fallback (copy your powerbb JSON to clipboard first):

python powerbb.py --output Q:\path\deck.pptx --template Q:\path\ExecDeck.pptx -v
# ➜ prompts you, shows a preview, and builds if you confirm
"""




from __future__ import annotations

import argparse
import json
import logging
import os
import re
import sys
from typing import Any, Dict, List, Optional, Tuple, Union
import platform
import subprocess

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

PowerBB = Dict[str, Any]

# Prefer the official location
try:
    from pptx.oxml.xmlchemy import OxmlElement  # python-pptx ≥ 0.6.x
except Exception:
    # Fallback: create an OxmlElement via lxml (python-pptx depends on lxml)
    from lxml import etree
    def OxmlElement(tag: str):
        """
        Accepts tags like 'a:buAutoNum' or Clark notation '{...}buAutoNum'
        and returns a namespaced element suitable for insertion.
        """
        if tag.startswith("{"):
            qname = tag
        else:
            qname = qn(tag)  # convert 'a:buAutoNum' -> '{ns}buAutoNum'
        return etree.Element(qname)

def _get_clipboard_text() -> Optional[str]:
    """
    Try several strategies to read plain-text from the system clipboard.
    Returns the text or None if unavailable.
    Order: pyperclip -> macOS pbpaste -> Linux xclip/xsel -> tkinter.
    """
    # 1) pyperclip (best cross-platform if installed)
    try:
        import pyperclip  # type: ignore
        txt = pyperclip.paste()
        if isinstance(txt, str) and txt:
            return txt
    except Exception:
        pass

    system = platform.system().lower()

    # 2) macOS pbpaste
    if "darwin" in system or "mac" in system:
        try:
            out = subprocess.check_output(["pbpaste"], stderr=subprocess.DEVNULL)
            if out:
                return out.decode("utf-8", errors="replace")
        except Exception:
            pass

    # 3) Linux xclip/xsel
    if "linux" in system:
        for cmd in (["xclip", "-selection", "clipboard", "-o"], ["xsel", "--clipboard", "--output"]):
            try:
                out = subprocess.check_output(cmd, stderr=subprocess.DEVNULL)
                if out:
                    return out.decode("utf-8", errors="replace")
            except Exception:
                continue

    # 4) tkinter fallback (works on Windows/mac/Linux with GUI)
    try:
        import tkinter as tk  # type: ignore
        r = tk.Tk()
        r.withdraw()
        try:
            txt = r.clipboard_get()
            return txt if isinstance(txt, str) and txt else None
        finally:
            r.destroy()
    except Exception:
        pass

    return None


# -----------------------------
# Logging
# -----------------------------
logger = logging.getLogger("powerbb")
handler = logging.StreamHandler()
handler.setFormatter(logging.Formatter("%(levelname)s: %(message)s"))
logger.addHandler(handler)
logger.setLevel(logging.INFO)

# -----------------------------
# Utilities
# -----------------------------

def _disable_bullets(paragraph):
    """Ensure paragraph shows no bullet/numbering from the template."""
    pPr = paragraph._element.get_or_add_pPr()
    # remove any existing list formatting
    for child in list(pPr):
        if child.tag in {qn('a:buChar'), qn('a:buAutoNum'), qn('a:buBlip'), qn('a:buNone')}:
            pPr.remove(child)
    pPr.append(OxmlElement('a:buNone'))

def _clear_list_props(paragraph) -> None:
    """Remove any existing bullet/numbering from a paragraph."""
    pPr = paragraph._element.get_or_add_pPr()
    for child in list(pPr):
        if child.tag in {qn('a:buChar'), qn('a:buAutoNum'), qn('a:buBlip'), qn('a:buNone')}:
            pPr.remove(child)

def _set_no_bullets(paragraph) -> None:
    """Force no bullets/numbering (a:buNone)."""
    pPr = paragraph._element.get_or_add_pPr()
    _clear_list_props(paragraph)
    pPr.append(OxmlElement('a:buNone'))

def _set_bullet(paragraph, char: str = "•") -> None:
    """Set a simple bullet character regardless of template defaults."""
    pPr = paragraph._element.get_or_add_pPr()
    _clear_list_props(paragraph)
    bu = OxmlElement('a:buChar')
    bu.set('char', char)
    pPr.append(bu)

def _set_numbering(paragraph, start_at: int | None = None, num_type: str = "arabicPeriod") -> None:
    """
    Use true PowerPoint numbering (a:buAutoNum).
    num_type examples: 'arabicPeriod', 'romanUcPeriod', 'alphaUcParen', etc.
    Only set startAt on the FIRST top-level paragraph of a numbered region.
    """
    pPr = paragraph._element.get_or_add_pPr()
    _clear_list_props(paragraph)
    bu = OxmlElement('a:buAutoNum')
    bu.set('type', num_type)
    if start_at is not None:
        bu.set('startAt', str(int(start_at)))
    pPr.append(bu)


_VAR_PATTERN = re.compile(r"\{\{([a-zA-Z0-9_\-]+)\}\}")

def _expand_vars(s: str, variables: Dict[str, str]) -> str:
    if not isinstance(s, str):
        return s
    def repl(m):
        key = m.group(1)
        return str(variables.get(key, m.group(0)))
    return _VAR_PATTERN.sub(repl, s)

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _ph_type_name(ph_type) -> str:
    # Map PP_PLACEHOLDER enum to readable names
    try:
        return PP_PLACEHOLDER(ph_type).name
    except Exception:
        return str(ph_type)

def _get_title_placeholder(slide):
    for shp in slide.shapes.placeholders:
        phf = getattr(shp, "placeholder_format", None)
        if phf and phf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return shp
    for shp in slide.shapes.placeholders:
        if getattr(shp, "has_text_frame", False):
            return shp
    return None

def _get_body_placeholders_sorted(slide):
    bodies = []
    for shp in slide.shapes.placeholders:
        phf = getattr(shp, "placeholder_format", None)
        if not phf:
            continue
        if phf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE):
            continue
        if getattr(shp, "has_text_frame", False):
            bodies.append(shp)
    bodies.sort(key=lambda s: s.left)
    return bodies

def _layout_body_count(slide_layout) -> int:
    cnt = 0
    for shp in slide_layout.placeholders:
        phf = getattr(shp, "placeholder_format", None)
        if not phf:
            continue
        if phf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE):
            continue
        cnt += 1
    return cnt

def _aspect_label(w, h) -> str:
    # Very rough aspect detector
    ratio = float(w) / float(h) if h else 0.0
    candidates = [(16/9, "16:9"), (4/3, "4:3"), (16/10, "16:10")]
    for target, label in candidates:
        if abs(ratio - target) < 0.02:
            return label
    return f"{ratio:.3f}:1"


def _build_template_profile(prs: Presentation) -> Dict[str, Any]:
    """
    Build a machine-readable profile of the template with:
    - slide size & aspect
    - masters, layouts, placeholder geometry
    - left/right indices for 2-body layouts
    - lists of two-body vs single-body layout names
    - alias suggestions and recommended defaults
    - a ready-to-paste powerbb 'meta' stub
    """
    profile: Dict[str, Any] = {
        "slide_size": {
            "width_emu": int(prs.slide_width),
            "height_emu": int(prs.slide_height),
            "aspect": _aspect_label(prs.slide_width, prs.slide_height),
        },
        "masters": []
    }
    for mi, m in enumerate(prs.slide_masters):
        m_info: Dict[str, Any] = {"name": m.name, "index": mi, "layouts": []}
        for li, l in enumerate(m.slide_layouts):
            phs = []
            for ph in l.placeholders:
                phf = getattr(ph, "placeholder_format", None)
                if not phf:
                    continue
                phs.append({
                    "idx": getattr(phf, "idx", None),
                    "type": _ph_type_name(phf.type),
                    "name": getattr(ph, "name", None),
                    "left": int(getattr(ph, "left", 0)),
                    "top": int(getattr(ph, "top", 0)),
                    "width": int(getattr(ph, "width", 0)),
                    "height": int(getattr(ph, "height", 0)),
                })
            body_slots = sum(1 for p in phs if p["type"] not in ("TITLE", "CENTER_TITLE", "SUBTITLE"))
            # Left/right indices by x ordering (only for 2+ bodies)
            body_phs = [p for p in phs if p["type"] not in ("TITLE", "CENTER_TITLE", "SUBTITLE")]
            body_phs_sorted = sorted(body_phs, key=lambda x: x["left"])
            lr = {"left_idx": None, "right_idx": None}
            if len(body_phs_sorted) >= 2:
                lr["left_idx"] = body_phs_sorted[0]["idx"]
                lr["right_idx"] = body_phs_sorted[1]["idx"]

            layout_info: Dict[str, Any] = {
                "index": li,
                "name": l.name,
                "body_slots": int(body_slots),
                "placeholders": phs,
                "left_right_hint": lr
            }
            m_info["layouts"].append(layout_info)
        profile["masters"].append(m_info)

    # Name buckets for convenience
    all_layout_names = [l["name"] for m in profile["masters"] for l in m["layouts"]]
    single_body = [l["name"] for m in profile["masters"] for l in m["layouts"] if l["body_slots"] == 1]
    two_body = [l["name"] for m in profile["masters"] for l in m["layouts"] if l["body_slots"] >= 2]

    profile["single_body_layouts"] = single_body
    profile["two_body_layouts"] = two_body
    profile["suggested_layout_aliases"] = _suggest_aliases(profile)
    profile["recommended_defaults"] = _recommended_defaults(profile)

    # Ready-to-paste powerbb meta stub (template_path to be filled by user)
    profile["powerbb_meta_stub"] = {
        "template_path": "<YOUR_TEMPLATE_PATH>.pptx",
        "default_layout": profile["recommended_defaults"].get("default_layout"),
        "layout_aliases": profile["suggested_layout_aliases"],
        "fallback_layout": profile["recommended_defaults"].get("fallback_layout"),
        "variables": {"client": "Acme", "year": "2025"},
        "defaults": {
            "list_type": "bullet",
            "fit": "shrink",
            "font_family": "Calibri",
            "title_size_pt": 40,
            "body_size_pt": 24
        }
    }
    return profile



def _suggest_aliases(profile: Dict[str, Any]) -> Dict[str, str]:
    """
    Heuristic mapping for common canonical names -> best matching layout in this template.
    """
    names = []
    for m in profile["masters"]:
        for l in m["layouts"]:
            names.append((l["name"], l["body_slots"]))

    def pick_two_col():
        # Prefer familiar names first, else max body_slots
        prefer = ["Two Content", "Comparison", "Title and Two Content"]
        for pref in prefer:
            for n, b in names:
                if n and pref.lower() in n.lower():
                    return n
        # else pick the first with >=2 bodies
        for n, b in names:
            if b >= 2:
                return n
        return None

    def pick_single():
        prefer = ["Title and Content", "Title, Content"]
        for pref in prefer:
            for n, b in names:
                if n and pref.lower() in n.lower():
                    return n
        # else any with 1 body
        for n, b in names:
            if b == 1:
                return n
        # last resort, any layout name
        return names[0][0] if names else None

    two = pick_two_col()
    one = pick_single()
    aliases = {}
    if two:
        aliases["two column with header"] = two
    if one:
        aliases["title + bullets"] = one
    return aliases

def _recommended_defaults(profile: Dict[str, Any]) -> Dict[str, str]:
    defaults = {"default_layout": None, "fallback_layout": None}
    # Try to use the single-content pick as default/fallback
    aliases = _suggest_aliases(profile)
    defaults["default_layout"] = aliases.get("title + bullets") or None
    defaults["fallback_layout"] = defaults["default_layout"]
    return defaults


def _dump_layouts(prs: Presentation, as_json: Optional[str] = None) -> None:
    profile = _build_template_profile(prs)

    ss = profile["slide_size"]
    logger.info(f"Slide size: {ss['width_emu']} x {ss['height_emu']} EMU  (~aspect {ss['aspect']})")
    for m in profile["masters"]:
        logger.info(f"[Master {m['index']}] {m.get('name') or ''}")
        for l in m["layouts"]:
            lr = l["left_right_hint"]
            lr_s = f" LR idx=({lr['left_idx']},{lr['right_idx']})" if lr["left_idx"] is not None else ""
            logger.info(f"  - [{m['index']}:{l['index']}] {l['name']} (body_slots={l['body_slots']}){lr_s}")
            for ph in l["placeholders"]:
                logger.info(
                    f"      ph idx={ph['idx']} type={ph['type']} name={ph['name']} "
                    f"pos=({ph['left']},{ph['top']}) size=({ph['width']},{ph['height']})"
                )

    logger.info("Two-body layouts: " + ", ".join(profile["two_body_layouts"]))
    logger.info("Single-body layouts: " + ", ".join(profile["single_body_layouts"]))
    logger.info(f"Suggested layout_aliases: {json.dumps(profile['suggested_layout_aliases'], ensure_ascii=False)}")
    logger.info(f"Recommended defaults: {json.dumps(profile['recommended_defaults'], ensure_ascii=False)}")
    logger.info("PowerBB meta stub: " + json.dumps(profile["powerbb_meta_stub"], ensure_ascii=False))

    if as_json:
        with open(as_json, "w", encoding="utf-8") as f:
            json.dump(profile, f, ensure_ascii=False, indent=2)
        logger.info(f"Wrote template profile JSON: {os.path.abspath(as_json)}")



def _apply_background(slide, bg_conf: Dict[str, Any]):
    if not bg_conf:
        return
    if "color" in bg_conf and bg_conf["color"]:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _hex_to_rgb(bg_conf["color"])


def _set_text_style(paragraph, style: Dict[str, Any],
                    font_family_default: Optional[str],
                    size_default_pt: Optional[int]):
    """
    Apply style to ALL runs in the paragraph (not just the first).
    This ensures template theme formatting doesn’t override a single-run edit.
    """
    style = style or {}
    if not paragraph.runs:
        paragraph.add_run()
    size_pt = style.get("size_pt", size_default_pt)

    for run in paragraph.runs:
        font = run.font
        if "bold" in style and style["bold"] is not None:
            font.bold = bool(style["bold"])
        if "italic" in style and style["italic"] is not None:
            font.italic = bool(style["italic"])
        if size_pt:
            font.size = Pt(float(size_pt))
        if "color" in style and style["color"]:
            font.color.rgb = _hex_to_rgb(style["color"])
        if font_family_default:
            font.name = font_family_default

def _enable_text_autofit(text_frame) -> None:
    _set_autofit_mode(text_frame, "shrink")
    try:
        text_frame.word_wrap = True
    except Exception:
        pass

def _set_autofit_mode(text_frame, mode: str) -> None:
    bodyPr = _get_or_add_bodyPr(text_frame)
    for el in list(bodyPr):
        if el.tag in {qn('a:noAutofit'), qn('a:spAutoFit'), qn('a:normAutofit')}:
            bodyPr.remove(el)
    tag = {'none': 'a:noAutofit', 'shrink': 'a:normAutofit', 'resize': 'a:spAutoFit'}[mode]
    bodyPr.append(OxmlElement(tag))



def _set_autofit_mode(text_frame, mode: str) -> None:
    """
    mode: 'none' | 'shrink' | 'resize'
      none   -> <a:noAutofit/>
      shrink -> <a:normAutofit/>
      resize -> <a:spAutoFit/>  (rarely used)
    Clears conflicting children first so toggles behave like the UI.
    """
    bodyPr = _get_or_add_bodyPr(text_frame)
    for el in list(bodyPr):
        if el.tag in {qn('a:noAutofit'), qn('a:spAutoFit'), qn('a:normAutofit')}:
            bodyPr.remove(el)
    tag = {'none': 'a:noAutofit', 'shrink': 'a:normAutofit', 'resize': 'a:spAutoFit'}[mode]
    bodyPr.append(OxmlElement(tag))

def _finalize_text_frame_autofit(shape, target_size_pt: Optional[int], font_family: Optional[str]) -> None:
    """
    Keep geometry fixed and make PowerPoint shrink text AND line spacing on overflow.
    Strategy:
      - Force <a:normAutofit lnSpcReduction="12000"/> (12% line-space reduction allowed)
      - Temporarily set auto_size=NONE, run fit_text() to downsize runs now
      - Restore auto_size=TEXT_TO_FIT_SHAPE so PPT keeps shrinking if user edits later
      - Ensure word_wrap
    """
    try:
        tf = shape.text_frame
    except Exception:
        return

    # ---- XML: <a:normAutofit lnSpcReduction="..."/> ----
    bodyPr = _get_or_add_bodyPr(tf)
    for el in list(bodyPr):
        if el.tag in {qn('a:noAutofit'), qn('a:spAutoFit'), qn('a:normAutofit')}:
            bodyPr.remove(el)
    norm = OxmlElement('a:normAutofit')
    # Allow some reduction of line spacing; value is in 1/1000 percent (100000 == 100%)
    # 12000 ≈ 12% reduction; tweak if you want more/less.
    norm.set('lnSpcReduction', '12000')
    bodyPr.append(norm)

    # ---- API: tighten wrap state ----
    try:
        tf.word_wrap = True
    except Exception:
        pass

    # ---- Run fit_text with auto_size temporarily NONE ----
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        prev_auto = getattr(tf, "auto_size", None)
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        prev_auto = None

    try:
        mx = int(target_size_pt) if target_size_pt else 24
        mn = max(9, int(mx * 0.55))  # allow more shrink; lots of bullets need it
        tf.fit_text(max_size=mx, max_lines=None, font_family=font_family, min_size=mn)
    except Exception:
        pass

    # ---- Restore shrink-on-overflow for future edits ----
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    _debug_autofit_state(shape, "[autofit]")


def _tighten_paragraph_spacing(p) -> None:
    """
    Remove template spacing that blocks autofit: zero space before/after and
    keep single-ish line spacing.
    """
    try:
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    except Exception:
        pass
    # 'line_spacing' accepts either a float multiplier or a Pt length.
    # Use a modest multiplier; leave a little breathing room.
    try:
        p.line_spacing = 1.0
    except Exception:
        pass




def _flatten_nodes(nodes: List[Dict[str, Any]],
                   variables: Dict[str, str] | None = None
                   ) -> List[Tuple[int, str, Dict[str, Any]]]:
    """
    Flatten BulletNodes to (level, text, style) without adding bullet/number text.
    """
    variables = variables or {}
    out: List[Tuple[int, str, Dict[str, Any]]] = []

    def walk(node_list, level: int):
        for node in node_list or []:
            text = _expand_vars(node.get("text", ""), variables)
            style = node.get("style", {}) or {}
            out.append((level, text, style))
            walk(node.get("children", []) or [], level + 1)

    walk(nodes or [], 0)
    return out


def _render_region(shape,
                   region: Dict[str, Any],
                   defaults: Dict[str, Any],
                   variables: Dict[str, str]) -> None:
    """
    Render bullets/numbering into the SHAPE's text_frame and then enforce shrink-to-fit.
    Writes only into the provided shape (which we choose to be the main text box).
    """
    tf = shape.text_frame
    list_type = (region.get("list_type") or defaults.get("list_type") or "bullet").lower()
    start_at = int(region.get("start_at") or 1)
    body_size_pt = defaults.get("body_size_pt")
    font_family = defaults.get("font_family")

    items = _flatten_nodes(region.get("bullets", []) or [], variables)

    # Clear once
    tf.clear()
    top_started = False

    for idx, (level, text, style) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        try:
            p.level = int(level)
        except Exception:
            p.level = 0

        if list_type == "number":
            _set_numbering(p, start_at=start_at if (level == 0 and not top_started) else None)
            if level == 0 and not top_started:
                top_started = True
        else:
            _set_bullet(p, char="•")

        _set_text_style(p, style, font_family, body_size_pt)

    _finalize_text_frame_autofit(shape, body_size_pt, font_family)


def _emu_to_in(v) -> float:
    try:
        return float(v) / 914400.0
    except Exception:
        return 0.0

def _log_slide_shapes(slide, slide_idx: int, layout_name: str, layout_id: str) -> None:
    logger.debug("[slide %d] layout [%s] %s", slide_idx, layout_id, layout_name)
    for shp in slide.shapes:
        phf = getattr(shp, "placeholder_format", None)
        ph_type = getattr(phf, "type", None)
        try:
            ph_name = ph_type.name if ph_type is not None else None
        except Exception:
            ph_name = str(ph_type)
        logger.debug(
            "  [shape] idx=%s type=%s name=%r pos=(%.2fin,%.2fin) size=(%.2fin×%.2fin) has_text=%s",
            getattr(phf, "idx", None),
            ph_name,
            getattr(shp, "name", None),
            _emu_to_in(getattr(shp, "left", 0)),
            _emu_to_in(getattr(shp, "top", 0)),
            _emu_to_in(getattr(shp, "width", 0)),
            _emu_to_in(getattr(shp, "height", 0)),
            getattr(shp, "has_text_frame", False),
        )

def _log_text_metrics(shape, label: str) -> None:
    try:
        tf = shape.text_frame
    except Exception:
        return
    sizes = []
    for p in tf.paragraphs:
        for r in p.runs:
            s = getattr(r.font, "size", None)
            if s is not None:
                try:
                    sizes.append(int(s.pt))
                except Exception:
                    pass
    bodyPr = _get_or_add_bodyPr(tf)
    flags = []
    for el in list(bodyPr):
        tag = re.sub(r"{.*}", "", el.tag)
        if tag in ("noAutofit", "spAutoFit", "normAutofit"):
            flags.append(tag)

    logger.debug(
        "%s pos=(%.2fin,%.2fin) size=(%.2fin×%.2fin) paras=%d sample_font_sizes=%s flags=%s",
        label,
        _emu_to_in(getattr(shape, "left", 0)),
        _emu_to_in(getattr(shape, "top", 0)),
        _emu_to_in(getattr(shape, "width", 0)),
        _emu_to_in(getattr(shape, "height", 0)),
        len(tf.paragraphs),
        sizes[:8],
        flags,
    )



def _fit_text_if_requested(text_frame, fit: Optional[str], defaults: Dict[str, Any]):
    if fit == "shrink":
        try:
            max_size = defaults.get("body_size_pt", 24)
            min_size = max(10, int(max_size * 0.6))
            family = defaults.get("font_family") or None
            text_frame.fit_text(max_size=max_size, max_lines=None,
                                font_family=family, min_size=min_size)
        except Exception:
            pass


def _resolve_layout(prs, slide_spec: dict,
                    template_path: str = None,
                    meta: dict = None,
                    logger=None,
                    strict: bool = False):
    """
    Resolve a python-pptx slide layout across ALL masters using any of:
      - slide_spec['layout']      (name; honors meta['layout_aliases'])
      - slide_spec['layout_id']   (token 'm:l' from dump_layouts)
      - slide_spec['like_slide']  (1-based slide number in the *template* file)
    Falls back to meta['default_layout'] → meta['fallback_layout'] → [0:0].
    Set strict=True to raise on not found / ambiguous name.
    """
    name = (slide_spec.get("layout") or "").strip()
    token = (slide_spec.get("layout_id") or "").strip()
    like_slide = slide_spec.get("like_slide")
    meta = meta or {}
    aliases = (meta.get("layout_aliases") or {})

    def log(msg):
        if logger: logger.info(msg)

    # Normalize name & apply alias
    if name:
        name = " ".join(name.split())  # collapse weird spacing
        if name in aliases:
            log(f"[resolver] alias '{name}' → '{aliases[name]}'")
            name = aliases[name]

    # Flat list of all layouts
    all_layouts = []
    for m_idx, master in enumerate(prs.slide_masters):
        for l_idx, lay in enumerate(master.slide_layouts):
            all_layouts.append((m_idx, l_idx, lay))

    # 1) Exact token m:l
    if token:
        try:
            m, l = map(int, token.split(":"))
            lay = prs.slide_masters[m].slide_layouts[l]
            log(f"[resolver] token '{token}' → [{m}:{l}] {lay.name}")
            return lay
        except Exception as e:
            if strict: raise ValueError(f"Bad layout_id '{token}': {e}")
            log(f"[resolver] WARN bad layout_id '{token}', trying other methods…")

    # 2) Name match across all masters (first match unless strict ambiguity)
    if name:
        matches = [(m, l, lay) for (m, l, lay) in all_layouts if (" ".join((lay.name or "").split()) == name)]
        if len(matches) == 1:
            m, l, lay = matches[0]
            log(f"[resolver] name '{name}' → [{m}:{l}] {lay.name}")
            return lay
        if len(matches) > 1:
            msg = f"Ambiguous layout name '{name}' matches: " + ", ".join([f"[{m}:{l}]" for (m,l,_) in matches])
            if strict: raise ValueError(msg)
            m, l, lay = matches[0]
            log(f"[resolver] WARN {msg}; using first [{m}:{l}]")
            return lay
        log(f"[resolver] name '{name}' not found; trying like_slide/defaults…")

    # 3) like_slide from template
    if like_slide and template_path:
        info = identify_slide_layout(template_path, int(like_slide))  # you already have this helper
        tok = info.get("layout_id")
        if tok:
            m, l = map(int, tok.split(":"))
            lay = prs.slide_masters[m].slide_layouts[l]
            log(f"[resolver] like_slide {like_slide} → [{m}:{l}] {lay.name}")
            return lay
        nm = info.get("layout_name")
        if nm:
            for (m, l, lay) in all_layouts:
                if lay.name == nm:
                    log(f"[resolver] like_slide name '{nm}' → [{m}:{l}]")
                    return lay

    # 4) Fallbacks
    for candidate in (meta.get("default_layout"), meta.get("fallback_layout")):
        if candidate:
            for (m, l, lay) in all_layouts:
                if " ".join((lay.name or "").split()) == " ".join(candidate.split()):
                    log(f"[resolver] fallback '{candidate}' → [{m}:{l}] {lay.name}")
                    return lay

    # Last resort
    m, l, lay = all_layouts[0]
    log(f"[resolver] FINAL fallback → [{m}:{l}] {lay.name}")
    return lay


def layout_token(prs, layout) -> str:
    for m, master in enumerate(prs.slide_masters):
        for l, lay in enumerate(master.slide_layouts):
            if lay is layout:
                return f"{m}:{l}"
    return "?:?"



def _extract_slide_info(slide):
    """Return dict with title, left/right (level,text) lists, and body slot count."""
    title_shp = _get_title_placeholder(slide)
    title_text = ""
    if title_shp and getattr(title_shp, "has_text_frame", False):
        title_text = "\n".join(p.text for p in title_shp.text_frame.paragraphs if p.text)
    bodies = _get_body_placeholders_sorted(slide)
    def extract_body(ph):
        if ph is None or not getattr(ph, "has_text_frame", False):
            return []
        return [(getattr(p, "level", 0) or 0, p.text) for p in ph.text_frame.paragraphs if p.text is not None]
    left = extract_body(bodies[0]) if len(bodies) >= 1 else []
    right = extract_body(bodies[1]) if len(bodies) >= 2 else []
    return {"title": title_text, "left": left, "right": right, "body_slots": len(bodies)}


def _find_slide_by_title(prs, expected_title: str):
    """Return _extract_slide_info(...) for the first slide whose title matches (normalized)."""
    want = _norm_text(expected_title).lower()
    for sl in prs.slides:
        info = _extract_slide_info(sl)
        if _norm_text(info["title"]).lower() == want:
            return info
    # Relax to 'contains' if exact not found (helps with templates adding prefixes)
    for sl in prs.slides:
        info = _extract_slide_info(sl)
        if want in _norm_text(info["title"]).lower():
            return info
    return None


# -----------------------------
# Core creation
# -----------------------------


def create_ppt_from_powerbb(
    powerbb: Union[str, PowerBB],
    output_path: str,
    template_path: Optional[str] = None,
) -> None:
    # Parse + normalize
    if isinstance(powerbb, str):
        pb: PowerBB = json.loads(powerbb)
    else:
        pb = powerbb
    pb = _prepare_powerbb(pb, normalize_escapes=True)

    meta = pb.get("meta", {}) or {}
    variables = meta.get("variables", {}) or {}
    defaults = meta.get("defaults", {}) or {}

    ppt_template = template_path or meta.get("template_path")
    prs = Presentation(ppt_template) if ppt_template else Presentation()

    if bool(meta.get("clear_existing")) and len(prs.slides) > 0:
        _remove_all_slides(prs)

    for i, slide_spec in enumerate(pb.get("slides", []), start=1):
        layout = _resolve_layout(
            prs, slide_spec=slide_spec, template_path=ppt_template,
            meta=meta, logger=logger, strict=False
        )
        logger.info("[build] slide %d: want name='%s', id='%s', like='%s' → using [%s] %s",
                    i, slide_spec.get("layout"), slide_spec.get("layout_id"),
                    slide_spec.get("like_slide"), layout_token(prs, layout), layout.name)
        slide = prs.slides.add_slide(layout)

        # Debug inventory
        _log_slide_shapes(slide, i, layout.name, layout_token(prs, layout))

        # Background
        _apply_background(slide, slide_spec.get("background") or {})

        # ----- Title -----
        title_txt = _expand_vars(slide_spec.get("title", "") or "", variables)
        title_ph = _get_title_placeholder(slide)
        if title_ph is not None and title_txt:
            tf = title_ph.text_frame
            _prime_text_frame_for_shrink(tf, defaults.get("title_size_pt"), defaults.get("font_family"))
            tf.clear()
            p = tf.paragraphs[0]  # after clear() there is one empty para
            p.text = title_txt
            _set_text_style(p, slide_spec.get("style", {}), defaults.get("font_family"), defaults.get("title_size_pt"))
            _finalize_text_frame_autofit(title_ph, defaults.get("title_size_pt"), defaults.get("font_family"))
            _log_textbox_metrics(title_ph, prefix="[title]")

        # ----- Body (main + optional secondary) -----
        main_ph, secondary_ph = _choose_main_and_secondary_text(slide)
        regs = slide_spec.get("regions", {}) or {}

        # LEFT region → main text box
        if regs.get("left") and main_ph is not None:
            tf = main_ph.text_frame
            _prime_text_frame_for_shrink(tf, defaults.get("body_size_pt"), defaults.get("font_family"))
            tf.clear()
            _log_textbox_metrics(main_ph, prefix="[before body]")
            _append_region_paragraphs(tf, regs["left"], defaults, variables, use_first_para=True)
            _finalize_text_frame_autofit(main_ph, defaults.get("body_size_pt"), defaults.get("font_family"))
            _log_textbox_metrics(main_ph, prefix="[after left]")

        # RIGHT region
        if regs.get("right") and main_ph is not None:
            # If a second suitable text box exists, use it; otherwise append to main.
            target = secondary_ph if secondary_ph is not None else main_ph
            tf = target.text_frame

            if target is secondary_ph:
                _prime_text_frame_for_shrink(tf, defaults.get("body_size_pt"), defaults.get("font_family"))
                tf.clear()
                _append_region_paragraphs(tf, regs["right"], defaults, variables, use_first_para=True)
            else:
                # Same box: add a spacer paragraph, then append (do NOT clear)
                spacer = tf.add_paragraph()
                _set_no_bullets(spacer)
                spacer.text = ""
                _append_region_paragraphs(tf, regs["right"], defaults, variables, use_first_para=False)

            _finalize_text_frame_autofit(target, defaults.get("body_size_pt"), defaults.get("font_family"))
            _log_textbox_metrics(target, prefix="[after right]")

        # Notes
        if slide_spec.get("notes"):
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            notes_tf.paragraphs[0].text = _expand_vars(slide_spec["notes"], variables)

    out_abs = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(out_abs), exist_ok=True)
    prs.save(out_abs)
    logger.info("Wrote: %s", out_abs)





def _prime_text_frame_for_shrink(tf, target_size_pt: Optional[int], font_family: Optional[str]) -> None:
    """
    Prepare a TextFrame so PowerPoint will shrink text on overflow.
    Safe to call before you add any text.
    """
    bodyPr = _get_or_add_bodyPr(tf)
    for el in list(bodyPr):
        if el.tag in {qn('a:noAutofit'), qn('a:spAutoFit'), qn('a:normAutofit')}:
            bodyPr.remove(el)
    bodyPr.append(OxmlElement('a:normAutofit'))
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        tf.word_wrap = True
    except Exception:
        pass

def _append_region_paragraphs(tf,
                              region: Dict[str, Any],
                              defaults: Dict[str, Any],
                              variables: Dict[str, str],
                              use_first_para: bool) -> None:
    """
    Append bullets/numbered paragraphs into an *existing* TextFrame.
    Also normalizes paragraph spacing so shrink can actually work.
    """
    list_type = (region.get("list_type") or defaults.get("list_type") or "bullet").lower()
    start_at = int(region.get("start_at") or 1)
    body_size_pt = defaults.get("body_size_pt")
    font_family = defaults.get("font_family")

    items = _flatten_nodes(region.get("bullets", []) or [], variables)

    first_used = False
    top_started = False

    for (level, text, style) in items:
        if use_first_para and not first_used and len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
            p.clear()
            first_used = True
        else:
            p = tf.add_paragraph()

        p.text = text
        try:
            p.level = int(level)
        except Exception:
            p.level = 0

        if list_type == "number":
            _set_numbering(p, start_at=start_at if (p.level == 0 and not top_started) else None)
            if p.level == 0 and not top_started:
                top_started = True
        else:
            _set_bullet(p, char="•")

        _set_text_style(p, style, font_family, body_size_pt)
        _tighten_paragraph_spacing(p)




def _get_or_add_bodyPr(text_frame):
    """
    Return the <a:bodyPr> element under this TextFrame's <a:txBody>,
    creating and inserting it as the first child if missing.
    Works on older python-pptx versions that lack get_or_add_bodyPr().
    """
    txBody = text_frame._element  # CT_TextBody
    # find existing bodyPr
    for child in txBody:
        if child.tag == qn('a:bodyPr'):
            return child
    # create and insert as first child (valid order: bodyPr, lstStyle?, p*)
    bodyPr = OxmlElement('a:bodyPr')
    if len(txBody):
        txBody.insert(0, bodyPr)
    else:
        txBody.append(bodyPr)
    return bodyPr



# -----------------------------
# Round-trip tests
# -----------------------------

def test_powerbb_roundtrip(tmp_output_path: Optional[str] = None, template_path: Optional[str] = None) -> None:
    out_path = tmp_output_path or os.path.join(os.getcwd(), "powerbb_roundtrip_test.pptx")
    powerbb_sample: PowerBB = {
        "meta": {
            "template_path": template_path or None,
            "default_layout": "Title and Content",
            "layout_aliases": {
                "two column with header": "Two Content",
                "Title + Bullets": "Title and Content"
            },
            "fallback_layout": "Title and Content",
            "variables": {"client": "TestClient", "year": "2025"},
            "defaults": {"list_type": "bullet", "fit": "shrink", "font_family": "Calibri",
                         "title_size_pt": 36, "body_size_pt": 20}
        },
        "slides": [
            {
                "layout": "two column with header",
                "title": "Executive Summary — {{client}} ({{year}})",
                "regions": {
                    "left": {
                        "list_type": "bullet",
                        "bullets": [
                            {"text": "Mission & context", "children": [
                                {"text": "Safety-critical systems"},
                                {"text": "Certification constraints"}
                            ]},
                            {"text": "Opportunities"}
                        ]
                    },
                    "right": {
                        "list_type": "number",
                        "start_at": 3,
                        "bullets": [
                            {"text": "Near-term wins"},
                            {"text": "12-month roadmap", "children": [
                                {"text": "Pilot → Scale → Institutionalize"}
                            ]}
                        ]
                    }
                },
                "notes": "Keep to 60 seconds."
            },
            {
                "layout": "Title and Content",
                "title": "Risks & Mitigations — {{client}}",
                "regions": {
                    "left": {
                        "bullets": [
                            {"text": "Model risk management", "children": [
                                {"text": "Data lineage & versioning"},
                                {"text": "Independent V&V"}
                            ]},
                            {"text": "Human factors", "children": [
                                {"text": "Procedural safeguards"},
                                {"text": "Training & adoption"}
                            ]}
                        ]
                    }
                },
                "background": {"color": "#FFFFFF"}
            }
        ]
    }

    # Build the deck with the user's template (which may already contain slides)
    create_ppt_from_powerbb(powerbb_sample, out_path, template_path=template_path)
    prs = Presentation(out_path)

    # Compute the exact expected titles after variable expansion
    expected_title_1 = "Executive Summary — TestClient (2025)"
    expected_title_2 = "Risks & Mitigations — TestClient"

    # Find our slides anywhere in the deck, by title
    s1 = _find_slide_by_title(prs, expected_title_1)
    s2 = _find_slide_by_title(prs, expected_title_2)

    if s1 is None or s2 is None:
        # Print all titles to aid debugging
        titles = [ _norm_text(_extract_slide_info(sl)["title"]) for sl in prs.slides ]
        raise AssertionError(
            f"Could not find test slides by title.\n"
            f"Expected: {_norm_text(expected_title_1)!r} and {_norm_text(expected_title_2)!r}\n"
            f"Found titles: {titles}"
        )

    # --- Diagnostics ---
    logger.info(f"[Diag] Found Slide 1 title: {s1['title']!r}")
    logger.info(f"[Diag] Slide 1 left: {s1['left']}")
    logger.info(f"[Diag] Slide 1 right: {s1['right']}")
    logger.info(f"[Diag] Slide 1 body_slots: {s1['body_slots']}")
    logger.info(f"[Diag] Found Slide 2 title: {s2['title']!r}")
    logger.info(f"[Diag] Slide 2 left: {s2['left']}")

    # --- Assertions for Slide 1 ---
    title1 = _norm_text(s1["title"])
    assert ("Executive Summary" in title1 and "TestClient (2025)" in title1), f"Title mismatch; got: {s1['title']!r}"

    expected_left = {
        (0, "Mission & context"),
        (1, "Safety-critical systems"),
        (1, "Certification constraints"),
        (0, "Opportunities"),
    }
    got_left = {(l, t) for (l, t) in s1["left"]}
    assert expected_left.issubset(got_left), f"Left bullets mismatch; missing {expected_left - got_left}; got {s1['left']}"

    # Numbering is a formatting property; just check text/levels exist (right may be merged)
    target = s1["right"] if s1["right"] else s1["left"]
    assert any((lvl == 0 and _norm_text(txt) == "Near-term wins") for (lvl, txt) in target), f"Missing 'Near-term wins'; got {target}"
    assert any((lvl == 0 and _norm_text(txt) == "12-month roadmap") for (lvl, txt) in target), f"Missing '12-month roadmap'; got {target}"
    assert any((lvl == 1 and _norm_text(txt) == "Pilot → Scale → Institutionalize") for (lvl, txt) in target), f"Missing nested item; got {target}"

    # --- Assertions for Slide 2 ---
    title2 = _norm_text(s2["title"])
    assert ("Risks & Mitigations" in title2 and "TestClient" in title2), f"Slide 2 title mismatch; got: {s2['title']!r}"
    expect2 = {
        "Model risk management", "Data lineage & versioning", "Independent V&V",
        "Human factors", "Procedural safeguards", "Training & adoption"
    }
    got2 = {t for (_, t) in s2["left"]}
    assert expect2.issubset(got2), f"Slide 2 left content mismatch; missing {expect2 - got2}; got {s2['left']}"

    # Notes check: locate the *notes* of the same slide we found as s1
    # Find the actual slide object corresponding to s1 by exact normalized title
    s1_slide = None
    want1 = _norm_text(expected_title_1).lower()
    for sl in prs.slides:
        if _norm_text(_extract_slide_info(sl)["title"]).lower() == want1:
            s1_slide = sl
            break
    if s1_slide:
        notes_text = getattr(s1_slide.notes_slide.notes_text_frame, "text", "")
        assert "Keep to 60 seconds" in _norm_text(notes_text), f"Notes mismatch; got: {notes_text!r}"

    logger.info(f"[OK] powerbb round-trip test passed. Output: {out_path}")


def _remove_all_slides(prs) -> None:
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)

def _body_text_placeholders(slide):
    """
    Return BODY-type text placeholders on the slide (exclude TITLE/SUBTITLE/PICTURE/OBJECT).
    Each item: (shape, area, name)
    """
    from pptx.enum.shapes import PP_PLACEHOLDER as PH
    items = []
    for shp in slide.shapes.placeholders:
        phf = getattr(shp, "placeholder_format", None)
        if not phf:
            continue
        ph_type = phf.type
        # Only text-bearing BODY-like placeholders
        if ph_type in (PH.TITLE, PH.CENTER_TITLE, PH.SUBTITLE, PH.OBJECT, PH.PICTURE):
            continue
        if not getattr(shp, "has_text_frame", False):
            continue
        area = int(getattr(shp, "width", 0)) * int(getattr(shp, "height", 0))
        items.append((shp, area, (getattr(shp, "name", "") or "")))
    # Largest first
    items.sort(key=lambda t: t[1], reverse=True)
    return items


def _choose_main_and_secondary_text(slide):
    """
    Return (main_text_shape, None). We deliberately ignore picture/obj placeholders
    and choose the largest text-capable BODY placeholder from slide.placeholders.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER as PH

    candidates = []
    for shp in slide.shapes.placeholders:
        phf = getattr(shp, "placeholder_format", None)
        if not phf:
            continue
        # Exclude non-body placeholders
        if phf.type in (PH.TITLE, PH.CENTER_TITLE, PH.SUBTITLE, PH.PICTURE, PH.OBJECT):
            continue
        if not getattr(shp, "has_text_frame", False):
            continue
        area = int(getattr(shp, "width", 0)) * int(getattr(shp, "height", 0))
        candidates.append((area, shp))

    # Fallback: any text frame that isn't a picture/object (rare in templates)
    if not candidates:
        for shp in slide.shapes:
            phf = getattr(shp, "placeholder_format", None)
            if phf and phf.type in (PH.PICTURE, PH.OBJECT, PH.TITLE, PH.CENTER_TITLE, PH.SUBTITLE):
                continue
            if getattr(shp, "has_text_frame", False):
                area = int(getattr(shp, "width", 0)) * int(getattr(shp, "height", 0))
                candidates.append((area, shp))

    candidates.sort(key=lambda t: t[0], reverse=True)
    main = candidates[0][1] if candidates else None

    # We intentionally do NOT return a secondary box; right-region content is appended to main.
    return (main, None)




def _emu_to_in(emu: int) -> float:
    return float(emu) / 914400.0

def _shape_dims_str(shp) -> str:
    try:
        w = _emu_to_in(int(shp.width)); h = _emu_to_in(int(shp.height))
        l = _emu_to_in(int(shp.left));  t = _emu_to_in(int(shp.top))
        return f"pos=({l:.2f}in,{t:.2f}in) size=({w:.2f}in×{h:.2f}in)"
    except Exception:
        return "pos/size=?"

def _log_textbox_metrics(shape, prefix: str = "") -> None:
    try:
        tf = shape.text_frame
    except Exception:
        logger.debug("%s [no text_frame] %s", prefix, _shape_dims_str(shape))
        return

    # Count paragraphs & runs, capture some font sizes
    para_count = len(tf.paragraphs)
    run_sizes = []
    for p in tf.paragraphs[:4]:
        for r in p.runs[:4]:
            sz = getattr(getattr(r, "font", None), "size", None)
            if sz is not None:
                try:
                    run_sizes.append(int(sz.pt))
                except Exception:
                    run_sizes.append(None)
    bodyPr = _get_or_add_bodyPr(tf)
    flags = []
    for el in list(bodyPr):
        t = re.sub(r"{.*}", "", el.tag)
        if t in ("noAutofit", "spAutoFit", "normAutofit"):
            flags.append(t)

    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        api_auto = getattr(tf, "auto_size", None)
        if api_auto is not None:
            # best-effort name
            for n in ("NONE", "TEXT_TO_FIT_SHAPE", "SHAPE_TO_FIT_TEXT"):
                if getattr(MSO_AUTO_SIZE, n, None) == api_auto:
                    api_auto = n
                    break
    except Exception:
        api_auto = getattr(tf, "auto_size", None)

    logger.debug(
        "%s %s paras=%d sample_font_sizes=%s flags=%s api_auto=%s wrap=%s",
        prefix, _shape_dims_str(shape), para_count, run_sizes[:6], flags, api_auto,
        getattr(tf, "word_wrap", None)
    )


def _fit_frame_text(text_frame, target_size_pt: Optional[int],
                    font_family: Optional[str], min_ratio: float = 0.6) -> None:
    """
    Conservative shrink-to-fit: enable autofit + call fit_text as a fallback.
    """
    _enable_text_autofit(text_frame)
    try:
        max_size = int(target_size_pt) if target_size_pt else 24
        min_size = max(10, int(max_size * float(min_ratio)))
        text_frame.fit_text(max_size=max_size, max_lines=None,
                            font_family=font_family, min_size=min_size)
    except Exception:
        # Some shapes don’t support fit_text; the XML autofit above still helps.
        pass


def _debug_autofit_state(shape, prefix: str = "") -> None:
    """Log the autofit flags present on a shape's text frame."""
    try:
        tf = shape.text_frame
    except Exception:
        return
    bodyPr = _get_or_add_bodyPr(tf)
    tags = []
    for el in list(bodyPr):
        t = re.sub(r"{.*}", "", el.tag)
        if t in ("noAutofit", "spAutoFit", "normAutofit"):
            tags.append(t)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        auto = getattr(tf, "auto_size", None)
        # convert to name if possible
        if auto is not None:
            for n in ("NONE", "TEXT_TO_FIT_SHAPE", "SHAPE_TO_FIT_TEXT"):
                if getattr(MSO_AUTO_SIZE, n, None) == auto:
                    auto = n
                    break
    except Exception:
        auto = getattr(tf, "auto_size", None)
    logger.debug("%s autofit=%s api_auto=%s wrap=%s", prefix, tags, auto, getattr(tf, "word_wrap", None))


def _build_powerbb_schema_text() -> str:
    """Returns a concise, copy-pasteable schema + usage notes for powerbb."""
    return (
        "You will produce one valid JSON object in the **powerbb** format (for python-pptx). "
        "Output ONLY JSON—no comments or prose.\n\n"
        "Schema (summary):\n"
        "{\n"
        '  "meta": {\n'
        '    "template_path": "path/to/template.pptx",   // required or provided externally\n'
        '    "default_layout": "Title and Content",\n'
        '    "layout_aliases": {"two column with header": "Two Content"},\n'
        '    "fallback_layout": "Title and Content",\n'
        '    "variables": {"client": "Acme", "year": "2025"},   // {{var}} interpolation\n'
        '    "defaults": { "list_type": "bullet", "fit": "shrink", "font_family": "Calibri", "title_size_pt": 40, "body_size_pt": 24 }\n'
        "  },\n"
        '  "slides": [\n'
        "    {\n"
        '      "layout": "Two Content",      // or alias; must exist in template\n'
        '      "master": "Executive Master", // optional master name\n'
        '      "title": "Slide Title — {{client}} ({{year}})",\n'
        '      "regions": {\n'
        '        "left":  { "list_type": "bullet", "bullets": [BulletNode, ...] },\n'
        '        "right": { "list_type": "number", "start_at": 1, "bullets": [BulletNode, ...] }\n'
        "      },\n"
        '      "notes": "Speaker notes.",\n'
        '      "tags": ["intro","exec"],\n'
        '      "style": {"title_color":"#0F2B5B","body_color":"#111111"},\n'
        '      "background": {"color":"#FFFFFF"}\n'
        "    }\n"
        "  ]\n"
        "}\n\n"
        "BulletNode (recursive): { \"text\": \"Point (supports {{variables}})\", \"style\": {\"bold\":false,\"italic\":false,\"color\":\"#111111\",\"size_pt\":24}, \"children\": [BulletNode, ...] }\n\n"
        "Usage notes:\n"
        "- Use ONLY real layout names from the provided template; if using synonyms, map via meta.layout_aliases.\n"
        "- Nesting depth maps to paragraph.level (0-based). Numbered lists use true PPT numbering; use start_at on the first top-level item.\n"
        "- Omit an entire region (left/right) if unused. Styling and fit are hints; template styles may override.\n"
        "- Provide concrete meta.variables so {{var}} placeholders resolve."
    )


def _build_template_specifics_text(profile: Dict[str, Any]) -> str:
    """Turn a template profile into human-readable guidance for the authoring prompt."""
    ss = profile["slide_size"]
    aliases = profile.get("suggested_layout_aliases", {})
    defaults = profile.get("recommended_defaults", {})
    meta_stub = profile.get("powerbb_meta_stub", {})

    two = ", ".join(profile.get("two_body_layouts", [])[:12]) or "(none found)"
    one = ", ".join(profile.get("single_body_layouts", [])[:12]) or "(none found)"

    return (
        f"Template specifics:\n"
        f"- Slide size (EMU): {ss['width_emu']} x {ss['height_emu']}  (~aspect {ss['aspect']})\n"
        f"- Two-body layouts (support left+right regions): {two}\n"
        f"- Single-body layouts (left only): {one}\n"
        f"- Suggested layout_aliases: {json.dumps(aliases, ensure_ascii=False)}\n"
        f"- Recommended defaults: {json.dumps(defaults, ensure_ascii=False)}\n\n"
        f"Ready-to-paste powerbb meta stub (fill template_path, adjust as needed):\n"
        "```json\n"
        f"{json.dumps(meta_stub, ensure_ascii=False, indent=2)}\n"
        "```\n"
        "Authoring guidance:\n"
        "- If you need two columns, choose a name from the two-body list.\n"
        "- If you need one column, use a name from the single-body list.\n"
        "- If you prefer friendly names ('two column with header'), add them under meta.layout_aliases mapping to real template names.\n"
    )


def generate_powerbb_prompt(template_path: Optional[str]) -> str:
    """
    Build a ready-to-paste prompt for a new conversation that explains powerbb and
    injects the target template's inventory so the LLM emits compatible JSON.
    """
    prs = Presentation(template_path) if template_path else Presentation()
    profile = _build_template_profile(prs)

    header = (
        "ROLE: You generate PowerPoint content as **powerbb JSON** for a Python builder.\n"
        "TASK: Produce ONE valid JSON object conforming to the powerbb schema.\n"
        "IMPORTANT: Output ONLY JSON—no prose, comments, or Markdown—unless explicitly asked otherwise.\n"
        "Output strict JSON only. Do not include Markdown escapes like \_, \&, \[ inside strings; write &, _, [, etc. directly. Do not wrap JSON in code fences.\n\n"
    )
    schema = _build_powerbb_schema_text() + "\n\n"
    templ = _build_template_specifics_text(profile)

    slide_content_guidance = """\n\nWhen asked to generate slide content, on each slide, aim for about 3-5 main 
    bullets and between 0-2 sub bullets under each, 
    typically with specific details, examples, or other supporting information from the research above. If no more 
    detailed information is needed, then it's fine for there to be no sub-bullets; use them only when they add value 
    to the slide and to the talk.\n\n"""


    return header + schema + templ + slide_content_guidance


# --- drop-in: put these near the top of your module ---
import json, os, re, logging

def clean_json_lenient(raw: str) -> str:
    """Make common non-JSON inputs parseable:
       - strip code fences, BOM
       - remove JS/C-style comments
       - unescape Markdown/Doc-ish backslashes like \[, \], \_, \&, \~
       - remove trailing commas before } or ]
    """
    s = raw

    # Strip UTF-8 BOM and surrounding whitespace
    s = s.lstrip("\ufeff").strip()

    # Strip Markdown/JSON code fences if present
    s = re.sub(r'^\s*```(?:json)?\s*', '', s, flags=re.IGNORECASE)
    s = re.sub(r'\s*```\s*$', '', s)

    # Remove // line comments and /* block comments */
    s = re.sub(r'(?m)^\s*//.*$', '', s)                # // comments
    s = re.sub(r'/\*.*?\*/', '', s, flags=re.DOTALL)   # /* ... */

    # Fix trailing commas: , } or , ]
    s = re.sub(r',\s*([\]}])', r'\1', s)

    # Remove backslashes that are NOT valid JSON escapes.
    # Valid escapes after \ are: " \ / b f n r t u
    s = re.sub(r'\\(?!["\\/bfnrtu])', '', s)

    # Normalize weird whitespace
    s = re.sub(r'\r\n?', '\n', s)

    return s


def load_deck_spec(path_or_text: str, lenient: bool, logger: logging.Logger):
    """Always read into a string, optionally clean, then json.loads."""
    if os.path.exists(path_or_text):
        with open(path_or_text, 'r', encoding='utf-8') as f:
            raw = f.read()
        logger.info(f"[input] loaded file '{path_or_text}' ({len(raw)} bytes)")
    else:
        raw = path_or_text
        logger.info(f"[input] using inline JSON ({len(raw)} chars)")

    cleaned = raw
    if lenient:
        cleaned = clean_json_lenient(raw)
        logger.info("[json] lenient cleanup applied")
    else:
        logger.info("[json] strict mode (no cleanup)")

    try:
        data = json.loads(cleaned)
        msg = "Parsed JSON successfully"
        if lenient:
            msg += " after lenient cleanup"
        logger.info(msg)
        return data
    except json.JSONDecodeError as e:
        logger.error(f"[json] parse failed at line {e.lineno} col {e.colno}: {e.msg}")
        # Drop a debug artifact to inspect exactly what was parsed
        try:
            with open("debug.cleaned.json", "w", encoding="utf-8") as dbg:
                dbg.write(cleaned)
            logger.info("[json] wrote cleaned candidate to debug.cleaned.json")
        except Exception as ioerr:
            logger.warning(f"[json] could not write debug.cleaned.json: {ioerr}")
        raise



def _strip_md_escapes(s: str) -> str:
    """Remove Markdown-style backslashes before punctuation inside an already-parsed Python string.
    Examples: 'V\\&V' -> 'V&V'; 'AI\\_SE' -> 'AI_SE'; 'list\\: item' -> 'list: item'
    Does NOT touch normal letters/numbers, so Windows paths 'C:\\temp\\file' remain unchanged.
    """
    import re
    if not isinstance(s, str) or not s:
        return s
    # Remove one-or-more backslashes that precede punctuation of interest.
    return re.sub(r"\\+([_\[\]\(\)\{\}&%#@!\+\-=:;,\.\?<>^~|/])", r"\1", s)




def _load_powerbb_from_file(path: str, lenient: bool = False) -> dict:
    """Read a file and parse JSON, optionally applying lenient cleanup."""
    with open(path, "r", encoding="utf-8") as f:
        txt = f.read()
    try:
        return json.loads(txt)
    except Exception as e:
        if not lenient:
            raise
        cleaned = clean_json_lenient(txt)
        return json.loads(cleaned)


def _normalize_powerbb_strings(obj):
    """Recursively normalize strings in PowerBB where they are human-facing text."""
    TEXT_KEYS = {"title", "text", "notes", "speaker_notes", "footer", "header", "alt_text"}
    if isinstance(obj, dict):
        out = {}
        for k, v in obj.items():
            if k in TEXT_KEYS and isinstance(v, str):
                out[k] = _strip_md_escapes(v)
            else:
                out[k] = _normalize_powerbb_strings(v)
        return out
    elif isinstance(obj, list):
        return [_normalize_powerbb_strings(x) for x in obj]
    else:
        return obj


def _prepare_powerbb(powerbb: dict, normalize_escapes: bool = True) -> dict:
    """Return a sanitized copy ready for rendering."""
    if normalize_escapes:
        return _normalize_powerbb_strings(powerbb)
    return powerbb


def _norm_text(s: str) -> str:
    """Normalize punctuation/whitespace so template quirks don't break tests."""
    import re
    s = (s or "")
    s = s.replace("\u2014", "-").replace("\u2013", "-").replace("—", "-").replace("–", "-")
    s = re.sub(r"\s+", " ", s).strip()
    return s


def cli_which_layout(template_path: str, slide_no: int):
    info = identify_slide_layout(template_path, slide_no)
    print(f"Slide {info['slide_number']} uses layout: {info['layout_name']}  (id [{info['layout_id']}])")


def identify_slide_layout(template_path: str, slide_number_1based: int):
    """
    Return dict with the layout name and [master:layout] token for a given slide number
    in the template deck. slide_number_1based is 1..N as shown in PowerPoint.
    """
    from pptx import Presentation
    prs = Presentation(template_path)
    if slide_number_1based < 1 or slide_number_1based > len(prs.slides):
        raise ValueError(f"Slide {slide_number_1based} out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number_1based - 1]
    target_layout = slide.slide_layout
    name = target_layout.name

    # Find [master:layout] indices
    master_idx = layout_idx = None
    for mi, master in enumerate(prs.slide_masters):
        for li, layout in enumerate(master.slide_layouts):
            if layout == target_layout:
                master_idx, layout_idx = mi, li
                break
        if master_idx is not None:
            break

    return {
        "slide_number": slide_number_1based,
        "layout_name": name,
        "layout_id": f"{master_idx}:{layout_idx}" if master_idx is not None else None,
    }




# -----------------------------
# CLI
# -----------------------------

def _parse_args(argv: List[str]) -> argparse.Namespace:
    import argparse
    p = argparse.ArgumentParser(
        description="Generate a PowerPoint deck from a powerbb JSON spec, inspect a template, generate a prompt, or run tests."
    )
    p.add_argument("--json", help="Path to powerbb JSON file. If omitted in build mode, you will be prompted to use clipboard content.")
    p.add_argument("--output", "-o", help="Output .pptx filepath (required unless --test / --dump-layouts / --promptgen).")
    p.add_argument("--template", "-t", help="Optional .pptx template providing masters/layouts.")
    p.add_argument("--test", action="store_true", help="Run the round-trip test suite and exit.")
    p.add_argument("--dump-layouts", action="store_true", help="Print a layout inventory for the chosen template/default.")
    p.add_argument("--dump-layouts-json", help="Write a machine-readable JSON template profile to this path.")
    p.add_argument("--promptgen", action="store_true", help="Generate a ready-to-paste authoring prompt (based on the provided --template).")
    p.add_argument("--prompt-out", help="Optional path to write the generated prompt text.")
    p.add_argument("--lenient", action="store_true", help="Attempt to auto-clean clipboard/file text before JSON parse.")
    p.add_argument("--verbose", "-v", action="count", default=0, help="Increase verbosity (-v, -vv).")
    
    return p.parse_args(argv)


def main(argv: List[str]) -> int:
    args = _parse_args(argv)

    if args.verbose >= 2:
        logger.setLevel(logging.DEBUG)
    elif args.verbose == 1:
        logger.setLevel(logging.INFO)
    else:
        logger.setLevel(logging.WARNING)

    # Open a presentation for optional dump/prompt
    try:
        prs_for_dump = Presentation(args.template) if args.template else Presentation()
    except Exception as e:
        logger.error(f"Unable to open template: {e}")
        return 1

    # Dump modes
    if args.dump_layouts or args.dump_layouts_json:
        _dump_layouts(prs_for_dump, as_json=args.dump_layouts_json)
        if not args.test and not args.json and not args.output and not args.promptgen:
            return 0

    # Prompt generation mode
    if args.promptgen:
        prompt_text = generate_powerbb_prompt(args.template)
        if args.prompt_out:
            with open(args.prompt_out, "w", encoding="utf-8") as f:
                f.write(prompt_text)
            logger.info(f"Wrote authoring prompt: {os.path.abspath(args.prompt_out)}")
        else:
            print(prompt_text)
        return 0

    # Test mode
    if args.test:
        try:
            test_powerbb_roundtrip(tmp_output_path=args.output, template_path=args.template)
            return 0
        except AssertionError as e:
            logger.error(f"Test failed: {e}")
            return 2
        except Exception as e:
            logger.error(f"Unexpected error during tests: {e}")
            return 3

    # ------- Build mode (interactive clipboard fallback for --json) -------
    if not args.output:
        logger.error("Missing --output. Provide an output .pptx path (or use --test / --dump-layouts / --promptgen).")
        return 1

    powerbb_obj = None

    if args.json:
        if not os.path.exists(args.json):
            logger.error(f"JSON file not found: {args.json}")
            return 1
        try:
            powerbb_obj = _load_powerbb_from_file(args.json, args.lenient)
            if args.lenient:
                logger.info("Parsed JSON after lenient cleanup.")

        except Exception as e:
        
            logger.error(f"Failed to read JSON: {e}")
            return 1
    else:
        if not sys.stdin.isatty():
            logger.error("No --json provided and not running interactively; cannot prompt for clipboard.")
            return 1

        print("No --json provided.")
        print("I can read the current clipboard and try to use it as your powerbb JSON.")
        resp = input("Use clipboard content as your powerbb JSON? [y/N]: ").strip().lower()
        if resp not in ("y", "yes"):
            logger.info("Aborted by user (no JSON source).")
            return 1

        clip = _get_clipboard_text()
        if not clip:
            logger.error("Could not read text from clipboard.")
            return 1

        # Show a short preview for confirmation
        preview = clip.strip().replace("\r\n", "\n")
        head = preview[:800]
        print("\n--- Clipboard preview (first 800 chars) ---\n" + head)
        if len(preview) > len(head):
            print("... [truncated]")
        resp2 = input("\nProceed with this clipboard content? [y/N]: ").strip().lower()
        if resp2 not in ("y", "yes"):
            logger.info("Aborted by user after preview; exiting.")
            return 1

        # Parse JSON
        try:
            powerbb_obj = json.loads(preview)
        except Exception as e:
            
            if args.lenient:
                cleaned = clean_json_lenient(preview)
                try:
                    powerbb_obj = json.loads(cleaned)
                    logger.info("Parsed clipboard JSON after lenient cleanup.")
                except Exception as e2:
                    logger.error(f"Clipboard content is not valid JSON (even after cleanup): {e2}")
                    return 1
            
            else:
                logger.error(f"Clipboard content is not valid JSON: {e}")
            
            return 1

    # Build
    try:
        create_ppt_from_powerbb(powerbb_obj, args.output, template_path=args.template)
        return 0
    except Exception as e:
        logger.error(f"Build failed: {e}")
        return 3




if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
