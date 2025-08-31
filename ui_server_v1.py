#!/usr/bin/env python3
# ui_server.py
from __future__ import annotations

import os
import sys
import json
import time
import tempfile
import threading
import subprocess
from dataclasses import dataclass
from typing import Optional, Dict, Any

from flask import Flask, request, jsonify, send_from_directory, Response
from flask import render_template_string

# ---------- CONFIG ----------
APP_NAME = "PowerBB UI"
APPDATA_DIR = os.path.join(os.environ.get("APPDATA", os.getcwd()), "PowerBB")
SETTINGS_PATH = os.path.join(APPDATA_DIR, "ui_settings.json")
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
POWERBB_PY = os.path.join(SCRIPT_DIR, "powerbb.py")  # assumes your script is here

# ---------- UTIL ----------
def ensure_dirs():
    os.makedirs(APPDATA_DIR, exist_ok=True)

def load_settings() -> Dict[str, Any]:
    ensure_dirs()
    if os.path.exists(SETTINGS_PATH):
        try:
            with open(SETTINGS_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    # defaults
    return {
        "template_path": "",
        "json_path": "",
        "output_path": "",
        "dump_json_path": "",
        "prompt_out": "",
        "lenient": True,
        "verbosity": 1,
        "open_output_after_build": True,
        "remember_settings": True,
        "json_source": "file",  # file | clipboard | inline
        "inline_json": ""
    }

def save_settings(s: Dict[str, Any]):
    ensure_dirs()
    with open(SETTINGS_PATH, "w", encoding="utf-8") as f:
        json.dump(s, f, ensure_ascii=False, indent=2)

import locale

def run_subprocess(cmd: list[str]) -> tuple[int, str]:
    enc = locale.getpreferredencoding(False) or "utf-8"
    proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    out = []
    while True:
        chunk = proc.stdout.readline()
        if not chunk and proc.poll() is not None:
            break
        if chunk:
            out.append(chunk.decode(enc, errors="replace"))
    code = proc.wait()
    return code, "".join(out)



# File dialogs via tkinter (local desktop)
def pick_file_open(title: str, filetypes: list[tuple[str, str]]) -> str | None:
    try:
        import tkinter as tk
        from tkinter import filedialog as fd
        root = tk.Tk()
        root.withdraw()
        # Force the dialog to the front and give it a parent
        root.attributes("-topmost", True)
        root.update_idletasks()
        try:
            path = fd.askopenfilename(title=title, filetypes=filetypes, parent=root)
        finally:
            root.destroy()
        return path or None
    except Exception:
        return None

def pick_file_save(title: str, defext: str, filetypes: list[tuple[str, str]], initialfile: str = "") -> str | None:
    try:
        import tkinter as tk
        from tkinter import filedialog as fd
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        root.update_idletasks()
        try:
            path = fd.asksaveasfilename(
                title=title,
                defaultextension=defext,
                filetypes=filetypes,
                initialfile=initialfile,
                parent=root,
            )
        finally:
            root.destroy()
        return path or None
    except Exception:
        return None


def get_clipboard_text() -> Optional[str]:
    # try pyperclip
    try:
        import pyperclip
        t = pyperclip.paste()
        if isinstance(t, str) and t.strip():
            return t
    except Exception:
        pass
    # PowerShell (Windows)
    try:
        import shutil
        ps = shutil.which("powershell") or shutil.which("pwsh")
        if ps and sys.platform.startswith("win"):
            out = subprocess.check_output([ps, "-NoProfile", "-Command", "Get-Clipboard"], stderr=subprocess.DEVNULL)
            return out.decode("utf-8", errors="replace")
    except Exception:
        pass
    # tkinter fallback
    try:
        import tkinter as tk
        r = tk.Tk(); r.withdraw()
        try:
            t = r.clipboard_get()
            return t if isinstance(t, str) and t.strip() else None
        finally:
            r.destroy()
    except Exception:
        pass
    return None

def write_temp_json(text: str) -> str:
    path = os.path.join(tempfile.gettempdir(), f"powerbb_ui_{int(time.time())}.json")
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)
    return path

# ---------- FLASK ----------
app = Flask(__name__, static_folder=None, template_folder=None)

# Serve the single-page app directly from string (we'll embed HTML below)

INDEX_HTML = r"""
<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>PowerBB AI Slide Manager</title>
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <script src="https://cdn.tailwindcss.com"></script>
  <style>
    .card { @apply bg-white rounded-2xl shadow p-6; }
    .btn  { @apply inline-flex items-center justify-center rounded-xl px-4 py-2 border border-gray-300 hover:bg-gray-50; }
    .btn-primary { @apply bg-orange-600 text-white border-orange-600 hover:bg-orange-700; }
    .btn-soft { @apply inline-flex items-center justify-center rounded-xl px-3 py-1.5 border border-gray-200 bg-gray-50 hover:bg-gray-100 text-xs; }
    .lbl { @apply text-sm text-gray-600; }
    .inp { @apply w-full rounded-xl border border-gray-300 px-3 py-2; }
    .row { @apply grid grid-cols-12 gap-4 items-center; }
    .pill { @apply inline-flex items-center px-2 py-1 rounded-full text-xs bg-gray-100 text-gray-700 mr-2; }
    .monos { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono","Courier New", monospace; }
    .tabbtn { @apply px-3 py-2 rounded-t-xl text-sm border border-b-0; }
    .tabbtn-active { @apply bg-white border-gray-300; }
    .tabbtn-inactive { @apply bg-gray-100 text-gray-600 border-transparent hover:bg-gray-200; }
    .disabled { @apply opacity-50 cursor-not-allowed; }
    .table { @apply w-full text-sm; }
    .th { @apply text-left text-gray-500 font-medium border-b border-gray-200 py-2; }
    .td { @apply border-b border-gray-100 py-2 align-top; }
    .idx { @apply text-gray-500 w-10; }
  </style>
</head>
<body class="bg-gray-100">
  <!-- Blue banner -->
  <div class="bg-gradient-to-r from-blue-900 to-blue-700 text-white">
    <div class="max-w-6xl mx-auto px-6 py-4">
      <h1 class="text-xl md:text-2xl font-semibold">PowerBB AI Slide Manager</h1>
      <div class="text-xs md:text-sm text-blue-200">Choose template → Inspect / Prompt → Provide JSON → Choose output → Generate slides</div>
    </div>
  </div>

  <div class="max-w-6xl mx-auto p-6 space-y-6">
    <!-- Step 1: Template -->
    <div class="card">
      <h2 class="text-lg font-semibold mb-4">1) Template</h2>
      <div class="row">
        <div class="col-span-10">
          <input id="template_path" size="120" class="inp monos" placeholder="(Optional) C:\\path\\to\\template.pptx — leave blank for None">
        </div>
        <div class="col-span-2 flex gap-2">
          <button class="btn w-full" onclick="pick('template_path','pptx','open','Select template')">Pick…</button>
          <button class="btn w-full" onclick="clearTemplate()">None</button>
        </div>
      </div>

      <div class="mt-4 grid md:grid-cols-3 gap-4">
        <!-- Dump layouts -->
        <div>
          <div class="flex items-center justify-between mb-2">
            <div class="font-medium">Dump layouts</div>
            <div class="flex gap-2">
              <button id="btnDump" class="btn-soft" onclick="runAction('dump_layouts')" title="Inspect the chosen template">Run</button>
              <button class="btn-soft" onclick="saveTextArea('dump_text','txt','template_layouts.txt')">Save As…</button>
            </div>
          </div>
          <textarea id="dump_text" class="inp monos h-40" placeholder="(Run after selecting a template)"></textarea>
        </div>

        <!-- Generate prompt -->
        <div>
          <div class="flex items-center justify-between mb-2">
            <div class="font-medium">Generate prompt</div>
            <div class="flex gap-2">
              <button id="btnPrompt" class="btn-soft" onclick="runAction('promptgen')" title="Prompt based on template">Run</button>
              <button class="btn-soft" onclick="saveTextArea('prompt_text','txt','prompt_for_llm.txt')">Save As…</button>
            </div>
          </div>
          <textarea id="prompt_text" class="inp monos h-40" placeholder="(Run to generate a ready-to-paste prompt)"></textarea>
        </div>

        <!-- NEW: List slides -->
        <div>
          <div class="flex items-center justify-between mb-2">
            <div class="font-medium">Slides in template</div>
            <div class="flex gap-2">
              <button id="btnListSlides" class="btn-soft" onclick="listSlides()">List slides</button>
              <button class="btn-soft" onclick="copyFrom('slides_csv_text')">Copy CSV</button>
              <button class="btn-soft" onclick="saveTextArea('slides_csv_text','csv','template_slides.csv')">Save As…</button>
            </div>
          </div>
          <div class="border rounded-xl overflow-hidden">
            <table class="table">
              <thead class="bg-gray-50">
                <tr>
                  <th class="th idx">#</th>
                  <th class="th w-24">Layout</th>
                  <th class="th">Layout name</th>
                  <th class="th">Title</th>
                </tr>
              </thead>
              <tbody id="slides_table_body"></tbody>
            </table>
          </div>
          <textarea id="slides_csv_text" class="hidden"></textarea>
        </div>
      </div>
    </div>

    <!-- Step 2: PowerBB JSON -->
    <div class="card">
      <h2 class="text-lg font-semibold mb-4">2) PowerBB JSON</h2>
      <div class="flex items-center gap-6">
        <label class="flex items-center gap-2"><input type="radio" name="jsonsrc" value="file" checked> File</label>
        <label class="flex items-center gap-2"><input type="radio" name="jsonsrc" value="inline"> Paste</label>
      </div>

      <div id="json_file_row" class="row mt-2">
        <div class="col-span-10">
          <input id="json_path" size="120" class="inp monos" placeholder="C:\\path\\to\\slides.powerbb.json">
        </div>
        <div class="col-span-2">
          <button class="btn w-full" onclick="pick('json_path','json','open','Select powerbb JSON')">Pick…</button>
        </div>
      </div>

      <div id="json_inline_row" class="mt-2 hidden">
        <textarea id="inline_json" class="inp monos h-56" placeholder='{"meta": {...}, "slides": [...]} (strict JSON)'></textarea>
        <div class="text-xs text-gray-500 mt-1">Tip: Paste raw JSON; use <b>Lenient JSON cleanup</b> in Settings if needed.</div>
      </div>
    </div>

    <!-- Step 3: Output -->
    <div class="card">
      <h2 class="text-lg font-semibold mb-4">3) Output deck (.pptx)</h2>
      <div class="row">
        <div class="col-span-10">
          <input id="output_path" size="120" class="inp monos" placeholder="C:\\path\\to\\deck.pptx">
        </div>
        <div class="col-span-2 flex gap-2">
          <button class="btn w-full" onclick="pick('output_path','pptx','open','Select existing .pptx')">Pick existing…</button>
          <button class="btn w-full" onclick="pickSave('output_path','pptx','Save deck as…','deck.pptx')">Save As…</button>
        </div>
      </div>
      <div class="text-xs text-gray-500 mt-2">Note: It’s fine if the template and output paths are the same file.</div>
    </div>

    <!-- Step 4: Generate -->
    <div class="card">
      <h2 class="text-lg font-semibold mb-4">4) Generate slides</h2>
      <div class="flex items-center gap-3">
        <button class="btn-primary" onclick="runAction('build')">Generate slides</button>
        <button class="btn" onclick="runAction('test')">Run self-test</button>
      </div>

      <div class="mt-6">
        <div class="lbl mb-1">CLI Equivalent <button class="btn-soft ml-2" onclick="copyFrom('cli')">Copy</button></div>
        <pre id="cli" class="bg-gray-900 text-gray-100 rounded-xl p-3 monos text-xs overflow-auto h-16"></pre>
      </div>

      <div class="border-b border-gray-300 flex gap-1 mt-4">
        <button id="tab-logs"   class="tabbtn tabbtn-active"   onclick="showTab('logs')">Logs</button>
        <button id="tab-prompt" class="tabbtn tabbtn-inactive" onclick="showTab('prompt')">Prompt</button>
        <button id="tab-dump"   class="tabbtn tabbtn-inactive" onclick="showTab('dump')">Dump</button>
      </div>

      <div id="panel-logs">
        <div class="flex items-center justify-between mt-3">
          <div class="lbl">Logs</div>
          <button class="btn-soft" onclick="copyFrom('logs')">Copy</button>
        </div>
        <pre id="logs" class="bg-black text-green-200 rounded-b-xl p-3 monos text-xs overflow-auto h-72"></pre>
      </div>

      <div id="panel-prompt" class="hidden">
        <div class="flex items-center justify-between mt-3">
          <div class="lbl">Generated Prompt</div>
          <div class="flex gap-2">
            <button class="btn-soft" onclick="copyFrom('prompt_text_tab')">Copy</button>
            <button class="btn-soft" onclick="saveTextArea('prompt_text_tab','txt','prompt_for_llm.txt')">Save As…</button>
          </div>
        </div>
        <textarea id="prompt_text_tab" class="inp monos h-72" placeholder="(Generate prompt in Step 1)"></textarea>
      </div>

      <div id="panel-dump" class="hidden">
        <div class="flex items-center justify-between mt-3">
          <div class="lbl">Dump Layouts</div>
          <div class="flex gap-2">
            <button class="btn-soft" onclick="copyFrom('dump_text_tab')">Copy</button>
            <button class="btn-soft" onclick="saveTextArea('dump_text_tab','txt','template_layouts.txt')">Save As…</button>
          </div>
        </div>
        <textarea id="dump_text_tab" class="inp monos h-72" placeholder="(Dump layouts in Step 1)"></textarea>
      </div>
    </div>

    <!-- Settings -->
    <div class="card">
      <h2 class="text-lg font-semibold mb-4">Settings</h2>
      <div class="grid md:grid-cols-2 gap-4">
        <label class="flex items-center gap-2"><input id="lenient" type="checkbox" checked> Lenient JSON cleanup</label>
        <label class="flex items-center gap-2"><input id="open_output_after_build" type="checkbox" checked> Open output folder after build</label>
      </div>
      <div class="grid md:grid-cols-2 gap-4 mt-3">
        <div>
          <label class="lbl">Verbosity</label>
          <select id="verbosity" class="inp">
            <option value="0">Normal</option>
            <option value="1" selected>-v</option>
            <option value="2">-vv</option>
          </select>
        </div>
        <div>
          <label class="lbl">Remember settings</label>
          <select id="remember_settings" class="inp">
            <option value="true" selected>Yes</option>
            <option value="false">No</option>
          </select>
        </div>
      </div>
    </div>

    <div class="text-xs text-gray-500">
      <span class="pill">Use Dump Layouts before authoring</span>
      <span class="pill">Enable Lenient for clipboard/Markdown JSON</span>
      <span class="pill">Settings persist in %APPDATA%/PowerBB</span>
    </div>
  </div>

<script>
const qs  = (s)=>document.querySelector(s);
const qsa = (s)=>Array.from(document.querySelectorAll(s));

async function api(path, opts={}) {
  const r = await fetch(path, Object.assign({headers:{'Content-Type':'application/json'}}, opts));
  if (!r.ok) throw new Error(await r.text());
  return await r.json();
}

function collectSettings() {
  return {
    template_path: qs('#template_path').value.trim(),
    json_path: qs('#json_path').value.trim(),
    output_path: qs('#output_path').value.trim(),
    lenient: qs('#lenient').checked,
    verbosity: parseInt(qs('#verbosity').value),
    open_output_after_build: qs('#open_output_after_build').checked,
    remember_settings: qs('#remember_settings').value === 'true',
    json_source: (qsa('input[name=jsonsrc]:checked')[0] || {}).value || 'file',
    inline_json: qs('#inline_json') ? qs('#inline_json').value : ""
  };
}
async function loadSettings() {
  try {
    const s = await api('/api/load-settings');
    qs('#template_path').value = s.template_path || '';
    qs('#json_path').value = s.json_path || '';
    qs('#output_path').value = s.output_path || '';
    qs('#lenient').checked = !!s.lenient;
    qs('#verbosity').value = String(s.verbosity ?? 1);
    qs('#open_output_after_build').checked = !!s.open_output_after_build;
    qs('#remember_settings').value = s.remember_settings ? 'true' : 'false';
    if (['file','inline'].includes(s.json_source)) {
      qsa('input[name=jsonsrc]').forEach(r => r.checked = (r.value === s.json_source));
    }
    if (qs('#inline_json')) qs('#inline_json').value = s.inline_json || '';
    toggleJsonSourceUI();
    updateTemplateActions();
  } catch(e) {}
}
async function saveSettings() {
  const s = collectSettings();
  if (!s.remember_settings) return;
  try { await api('/api/save-settings', {method:'POST', body:JSON.stringify(s)}); } catch(e) {}
}
qsa('input[name=jsonsrc]').forEach(r => r.addEventListener('change', ()=>{toggleJsonSourceUI(); saveSettings();}));
qsa('input,select,textarea').forEach(el => el.addEventListener('change', saveSettings));
function toggleJsonSourceUI(){
  const src = (qsa('input[name=jsonsrc]:checked')[0] || {}).value || 'file';
  qs('#json_file_row').classList.toggle('hidden', src !== 'file');
  qs('#json_inline_row').classList.toggle('hidden', src !== 'inline');
}
function updateTemplateActions(){
  const hasTpl = !!qs('#template_path').value.trim();
  ['btnDump','btnPrompt','btnListSlides'].forEach(id=>{
    const b = qs('#'+id);
    if (!b) return;
    b.classList.toggle('disabled', !hasTpl);
    b.disabled = !hasTpl;
  });
}
function clearTemplate(){
  qs('#template_path').value = "";
  updateTemplateActions();
  saveSettings();
}

async function pick(id, kind, mode, title) {
  const r = await api(`/api/pick-file?kind=${encodeURIComponent(kind)}&mode=${encodeURIComponent(mode)}&title=${encodeURIComponent(title)}&initial=${encodeURIComponent(qs('#'+id).value)}`);
  if (r.ok && r.path) { qs('#'+id).value = r.path; }
  else { alert(r.error || 'No file selected.'); }
  updateTemplateActions();
  saveSettings();
}
async function pickSave(id, kind, title, defname) {
  const r = await api(`/api/pick-file?kind=${encodeURIComponent(kind)}&mode=save&title=${encodeURIComponent(title)}&initial=${encodeURIComponent(defname)}`);
  if (r.ok && r.path) { qs('#'+id).value = r.path; saveSettings(); }
  else { alert(r.error || 'No file selected.'); }
}
function showTab(name){
  const names = ['logs','prompt','dump'];
  names.forEach(n=>{
    qs('#panel-'+n).classList.toggle('hidden', n!==name);
    qs('#tab-'+n).classList.toggle('tabbtn-active', n===name);
    qs('#tab-'+n).classList.toggle('tabbtn-inactive', n!==name);
  });
}
function setCLI(s){ qs('#cli').textContent = s || ''; }
function setLogs(s){ qs('#logs').textContent = s || ''; qs('#logs').scrollTop = qs('#logs').scrollHeight; }
function copyFrom(id){
  const el = qs('#'+id);
  const text = el.tagName === 'TEXTAREA' || el.tagName === 'INPUT' ? el.value : el.textContent;
  navigator.clipboard.writeText(text || '').catch(()=>{});
}

// NEW: list slides (template)
async function listSlides(){
  const path = qs('#template_path').value.trim();
  if (!path) { alert('Select a template first.'); return; }
  try{
    const res = await api('/api/list-slides', {method:'POST', body: JSON.stringify({path})});
    if (!res.ok) throw new Error(res.error || 'Unknown error');
    renderSlidesTable(res.slides || []);
    qs('#slides_csv_text').value = res.csv || '';
  }catch(e){
    alert('List failed.'); console.error(e);
  }
}
function renderSlidesTable(slides){
  const tb = qs('#slides_table_body');
  tb.innerHTML = '';
  slides.forEach(r=>{
    const tr = document.createElement('tr');
    tr.innerHTML = `
      <td class="td idx">${r.index}</td>
      <td class="td w-24"><code class="monos">${r.layout_id || ""}</code></td>
      <td class="td">${r.layout_name || ""}</td>
      <td class="td">${(r.title || "").replace(/</g,'&lt;').replace(/>/g,'&gt;')}</td>`;
    tb.appendChild(tr);
  });
}

async function runAction(action) {
  const s = collectSettings();
  const body = {
    action,
    template_path: s.template_path,
    json_source: s.json_source,
    json_path: s.json_path,
    inline_json: s.inline_json,
    output_path: s.output_path,
    dump_json_path: "",
    prompt_out: "",
    lenient: s.lenient,
    verbosity: s.verbosity,
    open_output_after_build: s.open_output_after_build
  };
  setCLI(''); setLogs('Running…');
  if (action==='promptgen') showTab('prompt');
  else if (action==='dump_layouts') showTab('dump');
  else showTab('logs');

  try {
    const res = await api('/api/run', {method:'POST', body:JSON.stringify(body)});
    setCLI(res.cli); setLogs(res.logs);
    if (res.prompt_text !== undefined) {
      const t = res.prompt_text || '';
      qs('#prompt_text').value = t;
      qs('#prompt_text_tab').value = t;
    }
    if (res.dump_text !== undefined) {
      const t = res.dump_text || '';
      qs('#dump_text').value = t;
      qs('#dump_text_tab').value = t;
    }
    if (!res.ok) { alert("Action failed (exit code " + res.exit_code + "). See Logs."); }
  } catch (e) {
    setLogs(String(e)); alert("Request failed—see Logs.");
  }
}

window.addEventListener('load', loadSettings);
</script>
</body>
</html>
"""




@app.get("/")
def index():
    return render_template_string(INDEX_HTML)

@app.get("/api/load-settings")
def api_load_settings():
    return jsonify(load_settings())

@app.post("/api/save-settings")
def api_save_settings():
    data = request.get_json(force=True)
    save_settings(data or {})
    return jsonify({"ok": True})

@app.get("/api/pick-file")
def api_pick_file():
    mode = request.args.get("mode", "open")
    kind = request.args.get("kind", "any")  # json|pptx|txt|any
    title = request.args.get("title", "Select file")
    initial = request.args.get("initial", "")

    ft_map = {
        "json": [("JSON files", "*.json"), ("All files", "*.*")],
        "pptx": [("PowerPoint", "*.pptx"), ("All files", "*.*")],
        "txt":  [("Text", "*.txt"), ("All files", "*.*")],
        "any":  [("All files", "*.*")]
    }
    filetypes = ft_map.get(kind, ft_map["any"])

    if mode == "save":
        path = pick_file_save(title, defext=f".{kind}" if kind in ("json","pptx","txt") else "", filetypes=filetypes, initialfile=os.path.basename(initial))
    else:
        path = pick_file_open(title, filetypes)

    return jsonify({"path": path})

@app.get("/api/get-clipboard")
def api_get_clipboard():
    txt = get_clipboard_text()
    return jsonify({"text": txt, "ok": bool(txt)})


def _read_text_if_exists(path: str) -> Optional[str]:
    """Read a text file if it exists; return None otherwise."""
    try:
        if path and os.path.exists(path):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception:
        pass
    return None


@app.post("/api/run")
def api_run():
    """
    JSON body:
    {
      "action": "build"|"test"|"dump_layouts"|"promptgen",
      "template_path": "...",
      "json_source": "file"|"clipboard"|"inline",
      "json_path": "...",
      "inline_json": "...",
      "output_path": "...",
      "dump_json_path": "...",
      "prompt_out": "...",
      "lenient": true/false,
      "verbosity": 0|1|2,
      "open_output_after_build": true/false
    }
    """
    data = request.get_json(force=True)
    action = (data.get("action") or "").strip()
    template_path = (data.get("template_path") or "").strip()
    json_source = (data.get("json_source") or "file").strip()
    json_path = (data.get("json_path") or "").strip()
    inline_json = (data.get("inline_json") or "")
    output_path = (data.get("output_path") or "").strip()
    dump_json_path = (data.get("dump_json_path") or "").strip()
    prompt_out = (data.get("prompt_out") or "").strip()
    lenient = bool(data.get("lenient", False))
    verbosity = int(data.get("verbosity", 0))
    open_after = bool(data.get("open_output_after_build", False))

    if not os.path.exists(POWERBB_PY):
        return jsonify({"ok": False, "logs": f"powerbb.py not found at {POWERBB_PY}"}), 400

    # Build command
    cmd = [sys.executable, POWERBB_PY]
    if template_path:
        cmd += ["--template", template_path]
    cmd += ["-v"] * max(0, min(2, verbosity))

    # Prepare per-action flags
    prompt_text = None
    dump_text = None
    dump_json_text = None

    if action == "build":
        # Determine JSON path (file | clipboard | inline)
        if json_source == "clipboard":
            clip = get_clipboard_text()
            if not clip:
                return jsonify({"ok": False, "logs": "Clipboard empty/unavailable."}), 400
            json_path = write_temp_json(clip)
        elif json_source == "inline":
            if not inline_json.strip():
                return jsonify({"ok": False, "logs": "No inline JSON provided."}), 400
            json_path = write_temp_json(inline_json)

        if not json_path:
            return jsonify({"ok": False, "logs": "JSON path is required for build."}), 400
        if not output_path:
            return jsonify({"ok": False, "logs": "Output path is required for build."}), 400

        cmd += ["--json", json_path, "--output", output_path]
        if lenient:
            cmd += ["--lenient"]

    elif action == "test":
        cmd += ["--test"]
        if output_path:
            cmd += ["--output", output_path]

    elif action == "dump_layouts":
        cmd += ["--dump-layouts"]
        if dump_json_path:
            cmd += ["--dump-layouts-json", dump_json_path]

    elif action == "promptgen":
        cmd += ["--promptgen"]
        # We *always* show the generated prompt inline;
        # also write to file if provided
        if prompt_out:
            cmd += ["--prompt-out", prompt_out]

    else:
        return jsonify({"ok": False, "logs": f"Unknown action: {action}"}), 400

    # Run
    code, logs = run_subprocess(cmd)

    # Collect inline artifacts
    if action == "promptgen":
        # powerbb prints the prompt to stdout; show inline always
        prompt_text = logs

        # If user also asked to save, the CLI already wrote a file; nothing more to do

    if action == "dump_layouts":
        # Human-readable dump is in stdout; show inline always
        dump_text = logs
        # If a JSON path was provided and created, include its content too (for copy/paste)
        dump_json_text = _read_text_if_exists(dump_json_path)

    result = {
        "ok": code == 0,
        "exit_code": code,
        "logs": logs,
        "cli": " ".join([f'"{c}"' if " " in c else c for c in cmd]),
        "prompt_text": prompt_text,
        "dump_text": dump_text,
        "dump_json_text": dump_json_text,
    }

    # Post-action nicety: open output folder on Windows
    if action == "build" and code == 0 and open_after and sys.platform.startswith("win"):
        try:
            folder = os.path.dirname(os.path.abspath(output_path))
            if folder and os.path.isdir(folder):
                os.startfile(folder)  # noqa
        except Exception:
            pass

    return jsonify(result)


@app.post("/api/save-text")
def api_save_text():
    """
    Save arbitrary text to a chosen path.
    Body: {"path":"C:\\file.txt", "text":"..."}
    """
    data = request.get_json(force=True)
    path = (data.get("path") or "").strip()
    text = data.get("text") or ""
    if not path:
        return jsonify({"ok": False, "error": "No path provided."}), 400
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(text)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


# --- run the Flask app locally ---
def _pick_free_port(preferred: int = 5000) -> int:
    import socket
    # try preferred first
    s = socket.socket()
    try:
        s.bind(("127.0.0.1", preferred))
        s.close()
        return preferred
    except OSError:
        s.close()
    # else pick any free port
    s2 = socket.socket()
    s2.bind(("127.0.0.1", 0))
    port = s2.getsockname()[1]
    s2.close()
    return port


@app.post("/api/list-slides")
def api_list_slides():
    """
    Body: {"path":"C:\\template.pptx"}
    Return: { ok, slides: [{index, layout_id, layout_name, title}], csv }
    """
    from pptx import Presentation
    try:
        data = request.get_json(force=True)
    except Exception:
        return jsonify({"ok": False, "error": "bad JSON"}), 400

    path = (data.get("path") or "").strip()
    if not path or not os.path.exists(path):
        return jsonify({"ok": False, "error": "file not found"}), 400

    prs = Presentation(path)

    # Build map from layout object -> (master_index, layout_index, name)
    layout_map = {}
    for mi, master in enumerate(prs.slide_masters):
        for li, layout in enumerate(master.slide_layouts):
            layout_map[id(layout)] = (mi, li, getattr(layout, "name", ""))

    # Helpers to extract a reasonable title
    def title_for_slide(slide):
        try:
            # Prefer TITLE or CENTER_TITLE placeholders
            for shp in slide.shapes:
                if getattr(shp, "is_placeholder", False):
                    pf = getattr(shp, "placeholder_format", None)
                    if pf and getattr(pf, "type", None) in (1, 3):  # 1=TITLE, 3=CENTER_TITLE
                        if getattr(shp, "has_text_frame", False):
                            t = shp.text_frame.text.strip()
                            if t:
                                return t
            # Fallback: any shape named like Title
            for shp in slide.shapes:
                nm = getattr(shp, "name", "").lower()
                if "title" in nm and getattr(shp, "has_text_frame", False):
                    t = shp.text_frame.text.strip()
                    if t:
                        return t
            # Fallback: first text
            for shp in slide.shapes:
                if getattr(shp, "has_text_frame", False):
                    t = shp.text_frame.text.strip()
                    if t:
                        return t.splitlines()[0]
        except Exception:
            pass
        return ""

    rows = []
    for i, sl in enumerate(prs.slides, start=1):
        layout = getattr(sl, "slide_layout", None)
        mi, li, lname = layout_map.get(id(layout), (None, None, getattr(layout, "name", "")))
        token = f"{mi}:{li}" if mi is not None else ""
        rows.append({
            "index": i,
            "layout_id": token,
            "layout_name": lname,
            "title": title_for_slide(sl),
        })

    # Compose CSV
    import csv, io
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(["#", "layout_id", "layout_name", "title"])
    for r in rows:
        w.writerow([r["index"], r["layout_id"], r["layout_name"], r["title"]])
    csv_text = buf.getvalue()

    return jsonify({"ok": True, "slides": rows, "csv": csv_text})




if __name__ == "__main__":
    ensure_dirs()
    port = _pick_free_port(5000)
    print(f"\n{APP_NAME} running at http://127.0.0.1:{port}\nPress Ctrl+C to stop.\n")
    # bind to loopback only; change host="0.0.0.0" if you need LAN access
    app.run(host="127.0.0.1", port=port, debug=False)

